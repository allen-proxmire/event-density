"""
Generate the 10-slide presentation deck for the ED 3-minute demo video.

Output: C:\\Users\\allen\\Desktop\\ED_3min_video_slides.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paths ──────────────────────────────────────────────────────────
# Script lives in outreach/slides/, so .. takes us to outreach/, ../.. to repo root.
REPO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..")
IMG_DIR = os.path.join(REPO, "outreach")
OUT_PATH = r"C:\Users\allen\Desktop\ED_3min_video_slides.pptx"

# ── Colors ─────────────────────────────────────────────────────────
BG        = RGBColor(0x12, 0x14, 0x20)   # deep navy
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC0, 0xC8, 0xD8)   # muted body
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # bright blue accent
DIM       = RGBColor(0x6B, 0x72, 0x80)   # dimmed text
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # checkmark green


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, text, left, top, width, height,
             font_size=24, color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=20, color=LIGHT, bold=False,
                  line_spacing=1.3, font_name="Calibri", align=PP_ALIGN.LEFT):
    """Add multiple lines as separate paragraphs."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, opts = line_info, {}
        else:
            text, opts = line_info

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(opts.get("size", font_size))
        p.font.color.rgb = opts.get("color", color)
        p.font.bold = opts.get("bold", bold)
        p.font.italic = opts.get("italic", False)
        p.font.name = opts.get("font", font_name)
        p.alignment = opts.get("align", align)
        p.space_after = Pt(int(font_size * (line_spacing - 1) + 2))

    return txBox


def add_image_safe(slide, path, left, top, width, height=None):
    """Add image if it exists, otherwise add a placeholder box."""
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top),
                                     Inches(width))
    else:
        add_text(slide, f"[Image: {os.path.basename(path)}]",
                 left, top, width, 0.5, font_size=14, color=DIM, italic=True)


# ══════════════════════════════════════════════════════════════════
#  BUILD DECK
# ══════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: Title ────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)

add_text(s, "Event Density",
         1.0, 1.8, 11.3, 1.5,
         font_size=60, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "A single PDE that reproduces known physics without importing it.",
         1.0, 3.5, 11.3, 1.2,
         font_size=28, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, "Allen Proxmire  |  github.com/Allen-Proxmire/Event-Density",
         1.0, 6.2, 11.3, 0.6,
         font_size=16, color=DIM, align=PP_ALIGN.CENTER)


# ── SLIDE 2: Four Primitives ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Four Primitives  \u2192  One PDE",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

primitives = [
    ("Density  \u03C1(x, t)", {"color": ACCENT, "bold": True, "size": 26}),
    ("    Bounded scalar field \u2014 the state variable", {"size": 22}),
    "",
    ("Mobility  M(\u03C1) = M\u2080(\u03C1max \u2212 \u03C1)\u1D5D", {"color": ACCENT, "bold": True, "size": 26}),
    ("    Degenerate diffusion \u2014 vanishes at capacity", {"size": 22}),
    "",
    ("Penalty  P(\u03C1) = P\u2080(\u03C1 \u2212 \u03C1*)", {"color": ACCENT, "bold": True, "size": 26}),
    ("    Monostable restoring force toward equilibrium", {"size": 22}),
    "",
    ("Participation  v(t)", {"color": ACCENT, "bold": True, "size": 26}),
    ("    Global feedback variable \u2014 non-local oscillation", {"size": 22}),
]
add_multiline(s, primitives, 1.2, 1.6, 10.5, 5.0, line_spacing=1.15)

add_text(s, "Seven axioms  \u2192  unique canonical PDE",
         0.8, 6.3, 11.7, 0.6,
         font_size=22, color=DIM, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 3: The Canonical PDE ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "The Canonical PDE",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

add_text(s, "\u2202\u209C\u03C1 = D [ \u2207\u00B7(M(\u03C1)\u2207\u03C1) \u2212 P(\u03C1) ] + Hv",
         1.0, 2.2, 11.3, 1.0,
         font_size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
         font_name="Consolas")

add_text(s, "v\u0307 = (F\u0304 \u2212 \u03B6v) / \u03C4",
         1.0, 3.5, 11.3, 0.8,
         font_size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
         font_name="Consolas")

channels = [
    ("Mobility", "Nonlinear diffusion  \u2192  Porous-medium equation"),
    ("Penalty", "Exponential relaxation  \u2192  RC circuit / Debye decay"),
    ("Participation", "Global oscillation  \u2192  Telegraph / RLC circuit"),
]
y = 4.8
for name, desc in channels:
    add_text(s, name, 2.0, y, 2.5, 0.5,
             font_size=22, color=ACCENT, bold=True)
    add_text(s, desc, 4.5, y, 7.0, 0.5,
             font_size=20, color=LIGHT)
    y += 0.55


# ── SLIDE 4: Running the Simulation ──────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Running the Simulation",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

specs = [
    ("2D grid:  64 \u00D7 64", {}),
    ("Parameters:  canonical defaults from axioms", {}),
    ("Runtime:  ~10 seconds", {}),
    ("Output:  density fields, energy, dissipation channels, 9 invariants", {}),
]
add_multiline(s, specs, 1.5, 2.0, 10.0, 3.0,
              font_size=26, line_spacing=1.6)

add_text(s, "[ Live notebook execution ]",
         1.0, 5.5, 11.3, 0.8,
         font_size=22, color=DIM, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 5: Density Field Evolution ─────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Density Field Evolution",
         0.8, 0.3, 11.7, 0.7,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

img_path = os.path.join(IMG_DIR, "plot_cell5_field_evolution.png")
add_image_safe(s, img_path, 0.8, 1.3, 11.7, 4.5)

add_text(s, "Cosine perturbation  \u2192  relaxation toward equilibrium \u03C1*",
         0.8, 6.2, 11.7, 0.6,
         font_size=20, color=DIM, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 6: RC Decay ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Penalty Channel  =  RC Decay",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

rc_text = [
    ("Silence mobility and participation (H = 0)", {}),
    ("", {}),
    ("PDE reduces to:  \u03B4\u0307 = \u2212DP\u2080\u03B4", {"color": ACCENT, "bold": True, "size": 24}),
    ("", {}),
    ("This IS the RC-circuit discharge equation.", {}),
    ("", {}),
    ("Error:  0.00%", {"color": GREEN, "bold": True, "size": 28}),
]
add_multiline(s, rc_text, 0.8, 1.6, 5.5, 4.5, font_size=22, line_spacing=1.2)

img_path = os.path.join(IMG_DIR, "plot_cell7_rc_decay.png")
add_image_safe(s, img_path, 6.8, 1.4, 5.8, 4.2)


# ── SLIDE 7: Telegraph ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Participation Channel  =  Telegraph Oscillation",
         0.8, 0.5, 11.7, 0.8,
         font_size=36, bold=True, align=PP_ALIGN.LEFT)

tel_text = [
    ("Turn participation on (H > 0)", {}),
    ("", {}),
    ("System becomes damped harmonic oscillator", {}),
    ("Identical to RLC circuit / telegraph equation", {}),
    ("", {}),
    ("Frequency error:   0.00%", {"color": GREEN, "bold": True, "size": 24}),
    ("Damping error:     0.09%", {"color": GREEN, "bold": True, "size": 24}),
]
add_multiline(s, tel_text, 0.8, 1.6, 5.5, 4.5, font_size=22, line_spacing=1.2)

img_path = os.path.join(IMG_DIR, "plot_cell8_telegraph.png")
add_image_safe(s, img_path, 6.8, 1.4, 5.8, 4.5)


# ── SLIDE 8: Nine Laws ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Nine Architectural Laws \u2014 Verified",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

laws = [
    "\u2713  1.  Unique attractor",
    "\u2713  2.  Monotone energy decay",
    "\u2713  3.  Spectral concentration",
    "\u2713  4.  Factorial complexity dilution",
    "\u2713  5.  Gradient-dissipation dominance",
    "\u2713  6.  Topological conservation",
    "\u2713  7.  Horizon formation",
    "\u2713  8.  Morphological hierarchy",
    "\u2713  9.  Sheet-filament oscillation",
]
law_items = [(law, {"color": GREEN, "size": 24}) for law in laws]
add_multiline(s, law_items, 2.0, 1.6, 9.0, 4.8, line_spacing=1.25)

add_text(s, "Hold across all parameter values and dimensions 1\u20134",
         0.8, 6.3, 11.7, 0.6,
         font_size=20, color=DIM, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 9: Beyond the Pipeline ─────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Beyond the Pipeline",
         0.8, 0.5, 11.7, 0.8,
         font_size=40, bold=True, align=PP_ALIGN.LEFT)

beyond = [
    ("10-material mobility law  (all R\u00B2 > 0.986)", {}),
    ("", {}),
    ("Dwarf galaxy rotation curves  \u2192  99.6% Burkert match", {}),
    ("", {}),
    ("Flat weak-lensing signal at 100\u20131000 kpc", {}),
    ("", {}),
    ("Baryonic Tully-Fisher exponent = 1/4  (observed)", {}),
    ("", {}),
    ("Two falsifiable predictions distinguishing ED from MOND and \u039BCDM", {}),
]
add_multiline(s, beyond, 1.5, 1.8, 10.0, 4.0, font_size=24, line_spacing=1.1)

add_text(s, "Everything is in the repository. Every claim has code behind it.",
         0.8, 6.0, 11.7, 0.8,
         font_size=22, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 10: Close ──────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, "Event Density",
         1.0, 1.5, 11.3, 1.2,
         font_size=60, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Clone it.  Run it.  Tell me what you see.",
         1.0, 3.2, 11.3, 1.0,
         font_size=32, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "github.com/Allen-Proxmire/Event-Density",
         1.0, 4.6, 11.3, 0.7,
         font_size=24, color=LIGHT, align=PP_ALIGN.CENTER,
         font_name="Consolas")

add_text(s, "Open in Colab  \u2192  outreach/ED_minimal_demo.ipynb",
         1.0, 5.3, 11.3, 0.6,
         font_size=20, color=DIM, align=PP_ALIGN.CENTER)

add_text(s, "Allen Proxmire  |  2026",
         1.0, 6.5, 11.3, 0.5,
         font_size=16, color=DIM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Size:  {os.path.getsize(OUT_PATH) / 1024:.0f} KB")
print(f"Slides: {len(prs.slides)}")
